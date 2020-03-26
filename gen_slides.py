from pptx import Presentation
from pptx.util import Inches, Pt
from os import system
from sys import platform
import os
import requests
import argparse
import re
import urllib
import shutil
from pptx.enum.text import MSO_AUTO_SIZE
from textwrap import wrap

# =================== INIT ARGS
if platform == "win32":
    system("color")
parser = argparse.ArgumentParser(
    description="Generates a history presentation.")
parser.add_argument("term", type=str, help="term to search")
parser.add_argument(
    "--search", "-s", help="additional search parameters", default="")
parser.add_argument(
    "--author", "-a", help="add the author of presentation", default="")
parser.add_argument(
    "--no-best", help="exclude the best defenition slide", action="store_true", default=False)
parser.add_argument(
    "--no-sources", help="exclude sources slide", action="store_true", default=False)

args = parser.parse_args()
term = args.term
params = args.search
presentation_params = {
    "author": args.author,
    "best": not args.no_best,
    "sources": not args.no_sources
}


def find_defenitions(term, params):
    # =================== FIND PAGES
    if not params:
        params = ""
    url = "https://duckduckgo.com/html/"

    print("Requesting url: {}...".format(url))

    headers = {
        "user-agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Brave Chrome/80.0.3987.149 Safari/537.36"
    }
    search_results = requests.post(url, {
        "kl": "ru-ru",
        "q": term + " " + params
    }, headers=headers).text

    print(term + params)

    url_regexp = "http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\(\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+"

    pattern = "class=\"result__a\"\\s*href=\"\\s*"
    pattern += url_regexp
    pattern += "\\s*\">"

    hrefs = re.findall(pattern, search_results)
    print("Found {} results!".format(len(hrefs)))

    cleanr = re.compile("<.*?>|&([a-z0-9]+|#[0-9]{1,6}|#x[0-9a-f]{1,6});")
    bodyonly = re.compile("^(.|\s)*<\s*body\s*>|<\s*/\s*body\s*>(.|\s)*$")
    defenition_regexp = re.compile(
        "({}\s*(((—|-|,)\s*(это))|(—|-|,)|(это))\s*([а-яёА-ЯЁ ,0-9/;]|([.]\s+[^А-Я0-9]|[.][а-я]))+?([.]|\n))\s+([\nА-Я0-9]|$)".format(term), re.IGNORECASE)

    print("="*14+"DEFENITIONS="+"="*14)

    # =================== PARSE DEFENITIONS
    result_defenitions = {}
    for href in hrefs:
        link = re.findall(url_regexp, href)[0]
        domain = re.findall("[^/:]+/?", link)[1].strip("/")
        print("\033[33mScanning \033[34m{}\033[33m...\033[0m".format(
            domain) + " "*25, end="\r")

        try:
            site_html = requests.get(link, headers=headers).text
        except KeyboardInterrupt:
            print("Exiting..." + " "*40)
            exit()
        except:
            continue

        site_html = re.sub("</(p|h[0-9]|)>", ". ", site_html)
        body_from = site_html.find("<body>") + 6
        body_to = site_html.find("</body>")
        site_text = re.sub(cleanr, "", site_html[body_from:body_to])
        site_text = re.sub("[(][^)]+[)]|[)]", "", site_text)

        defenitions = re.findall(defenition_regexp, site_text)
        defenitions = filter(lambda x: x[1] != "," or re.match(
            "[А-ЯЁA-Z]", x[0][1]), defenitions)
        defenitions = [re.sub("\s+", " ", x[0]) for x in defenitions]

        if (defenitions):
            defenition = max(defenitions, key=len)
            defenition = re.sub("\s+([.,])", "\g<1>", defenition)
            if (len(defenition.split(" ")) > 6):
                if (defenition not in result_defenitions.values()):
                    result_defenitions[link] = defenition
                    print(defenition)
                    print("="*40)

    return result_defenitions


def find_images(term):
    # =================== REQUEST IMAGES
    headers = {
        "user-agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Brave Chrome/80.0.3987.149 Safari/537.36"
    }
    url = "https://yandex.ru/images/search?text={}".format(term)

    print("Requesting url: {}...".format(url))
    api_result = requests.post(url, headers=headers).text

    image_regexp = "https://avatars.mds.yandex.net/get-pdb/(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\(\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+/"
    image_links = [x+"orig" for x in set(re.findall(image_regexp, api_result))]

    print("Found {} images!".format(len(image_links)))
    print("Downloading...", end="\r")

    # =================== DOWNLOAD IMAGES
    if not os.path.exists(term + "_images"):
        os.makedirs(term + "_images")

    images = []
    for i, image_link in enumerate(image_links):
        print("Downloading... [{}/{}]".format(i+1, len(image_links)), end="\r")
        path = "{}_images/{}.jpg".format(term, i)
        urllib.request.urlretrieve(image_link, path)
        images.append(path)

    print("Successfully downloaded!")
    return(images)


def create_presentation(defenitions, images, params):
    print("Generating slides...")
    presentation = Presentation()
    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = term.title()
    slide.placeholders[1].text = params["author"]

    for source, defenition in defenitions.items():
        layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(layout)

        domain = re.findall("[^/:]+/?", source)[1].strip("/")
        slide.shapes.title.text = domain
        slide.shapes.placeholders[1].text_frame.text = defenition

    if (params["best"]):
        layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(layout)

        best_defenition = max(defenitions.values(), key=len)
        slide.shapes.title.text = "Лучшее Определение"
        slide.shapes.placeholders[1].text_frame.text = best_defenition

    if (params["sources"]):
        layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = "Источники"
        frame = slide.shapes.placeholders[1].text_frame
        frame.text = list(defenitions.keys())[0]
        for source in list(defenitions.keys())[1:]:
            paragraph = frame.add_paragraph()
            paragraph.level = 0
            paragraph.text = source

        for paragraph in frame.paragraphs:
            paragraph.font.size = Pt(12)

    for image in images:
        layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(layout)

        slide.shapes.add_picture(image, Inches(
            1.75), Inches(1.25), height=Inches(5))

    presentation.save("{}.pptx".format(term))


# =================== RUN SCRIPT
print("Finding defenitions...")
print("="*40)
defenitions = find_defenitions(term, params)
print("Finding images..." + " "*20)
print("="*40)
images = find_images(term)
print("Creating presentation...")
print("="*40)
create_presentation(defenitions, images, presentation_params)
print("Cleaning up...")
if os.path.exists("{}_images".format(term)):
    shutil.rmtree("{}_images".format(term))
